"""
Data processors that convert raw data into text.
These processors use tools_engine to extract text from various data formats.
"""
import io
import logging
import tempfile
import os
from typing import Dict, Any, Optional

from .interfaces import DataProcessor
from tools_engine.client import get_tools_client

logger = logging.getLogger(__name__)


class BaseProcessor(DataProcessor):
    """Base class for all processors."""
    
    def __init__(self):
        """Initialize the processor with tools client."""
        self.tools_client = get_tools_client()
    
    def process(self, data: bytes, metadata: Dict[str, Any]) -> str:
        """
        Process binary data into text.
        
        Args:
            data: Raw binary data
            metadata: Metadata about the data source
            
        Returns:
            Processed text
        """
        # Default implementation just decodes as UTF-8
        try:
            return data.decode('utf-8', errors='replace')
        except Exception as e:
            logger.error(f"Error processing data: {str(e)}")
            return ""


class PDFProcessor(BaseProcessor):
    """Processor for PDF files."""
    
    def process(self, data: bytes, metadata: Dict[str, Any]) -> str:
        """
        Process PDF data into text.
        
        Args:
            data: Raw PDF data
            metadata: Metadata about the PDF
            
        Returns:
            Extracted text from the PDF
        """
        try:
            # Use tools_client to convert PDF to text
            result = self.tools_client.convert_pdf_to_text(
                pdf_data=data,
                extract_tables=True
            )
            
            if not result.get('success', False):
                logger.error(f"Error extracting text from PDF: {result.get('error_message')}")
                return ""
            
            # Combine extracted text and table data
            text = result.get('extracted_text', '')
            table_data = result.get('table_data', [])
            
            # Format table data as text
            table_text = ""
            for i, table in enumerate(table_data):
                table_text += f"\n\nTable {i+1}:\n{table}\n"
            
            return text + table_text
            
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            return ""


class ExcelProcessor(BaseProcessor):
    """Processor for Excel and CSV files."""
    
    def process(self, data: bytes, metadata: Dict[str, Any]) -> str:
        """
        Process Excel/CSV data into text.
        
        Args:
            data: Raw Excel/CSV data
            metadata: Metadata about the file
            
        Returns:
            Structured text representation of the spreadsheet
        """
        try:
            # Create a temporary file
            extension = metadata.get('file_extension', '.xlsx')
            with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as temp_file:
                temp_file.write(data)
                temp_file_path = temp_file.name
            
            try:
                # For CSV, we can use a simple format conversion
                if extension == '.csv':
                    # Convert CSV to JSON for easier processing
                    result = self.tools_client.convert_data_format(
                        input_data=data.decode('utf-8', errors='replace'),
                        input_format='csv',
                        output_format='json'
                    )
                    
                    if not result.get('success', False):
                        logger.error(f"Error converting CSV: {result.get('error_message')}")
                        return ""
                    
                    # Format the JSON as readable text
                    return f"CSV Data:\n{result.get('converted_data', '')}"
                else:
                    # For Excel, we need custom processing
                    # In a real implementation, you would use Python code here to process Excel
                    script = """
import pandas as pd
import json
import sys

# Read the Excel file
file_path = sys.argv[1]
dfs = {}

# Read all sheets
excel_file = pd.ExcelFile(file_path)
for sheet_name in excel_file.sheet_names:
    df = pd.read_excel(file_path, sheet_name=sheet_name)
    dfs[sheet_name] = df.to_json(orient='records')

# Output as JSON
print(json.dumps(dfs))
"""
                    
                    # Execute the script
                    result = self.tools_client.execute_python_script(
                        script_content=script,
                        parameters={"file_path": temp_file_path}
                    )
                    
                    if not result.get('success', False):
                        logger.error(f"Error processing Excel: {result.get('error_message')}")
                        return ""
                    
                    # Format the output as readable text
                    output_data = result.get('output_data', '{}')
                    return f"Excel Data:\n{output_data}"
            finally:
                # Clean up the temporary file
                if os.path.exists(temp_file_path):
                    os.remove(temp_file_path)
                    
        except Exception as e:
            logger.error(f"Error processing Excel/CSV: {str(e)}")
            return ""


class PowerPointProcessor(BaseProcessor):
    """Processor for PowerPoint files."""
    
    def process(self, data: bytes, metadata: Dict[str, Any]) -> str:
        """
        Process PowerPoint data into text.
        
        Args:
            data: Raw PowerPoint data
            metadata: Metadata about the file
            
        Returns:
            Extracted text from the presentation
        """
        try:
            # Create a temporary file
            extension = metadata.get('file_extension', '.pptx')
            with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as temp_file:
                temp_file.write(data)
                temp_file_path = temp_file.name
            
            try:
                # Extract text from PowerPoint using Python script
                script = """
import python_pptx
from pptx import Presentation
import json
import sys

def extract_text_from_slide(slide):
    text = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text.append(shape.text)
    return text

file_path = sys.argv[1]
ppt = Presentation(file_path)

slides = []
for i, slide in enumerate(ppt.slides):
    slide_text = extract_text_from_slide(slide)
    slides.append({
        "slide_number": i+1,
        "content": "\\n".join(slide_text)
    })

print(json.dumps({"slides": slides}))
"""
                
                # Execute the script
                result = self.tools_client.execute_python_script(
                    script_content=script,
                    parameters={"file_path": temp_file_path}
                )
                
                if not result.get('success', False):
                    logger.error(f"Error processing PowerPoint: {result.get('error_message')}")
                    return ""
                
                # Format the output as readable text
                output_data = result.get('output_data', '{"slides": []}')
                try:
                    import json
                    slides_data = json.loads(output_data).get('slides', [])
                    
                    text = "PowerPoint Presentation:\n\n"
                    for slide in slides_data:
                        text += f"Slide {slide.get('slide_number')}:\n"
                        text += f"{slide.get('content')}\n\n"
                    
                    return text
                except:
                    return f"Presentation Data:\n{output_data}"
            finally:
                # Clean up the temporary file
                if os.path.exists(temp_file_path):
                    os.remove(temp_file_path)
                    
        except Exception as e:
            logger.error(f"Error processing PowerPoint: {str(e)}")
            return ""


class AudioProcessor(BaseProcessor):
    """Processor for audio files using speech-to-text."""
    
    def process(self, data: bytes, metadata: Dict[str, Any]) -> str:
        """
        Process audio data into text using speech-to-text.
        
        Args:
            data: Raw audio data
            metadata: Metadata about the file
            
        Returns:
            Transcribed text from the audio
        """
        try:
            # In a real implementation, you would use Google Cloud Speech-to-Text API
            # via the tools_engine client
            
            # This is a placeholder for the API call
            logger.info(f"Processing audio file: {metadata.get('filename')}")
            
            # Create a temporary file
            extension = metadata.get('file_extension', '.mp3')
            with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as temp_file:
                temp_file.write(data)
                temp_file_path = temp_file.name
            
            try:
                # Example Python script to call Google Cloud Speech-to-Text API
                # In a real implementation, this would be a proper API call
                script = """
import io
import os
from google.cloud import speech

# Initialize client
client = speech.SpeechClient()

# Read the audio file
with io.open(sys.argv[1], "rb") as audio_file:
    content = audio_file.read()

audio = speech.RecognitionAudio(content=content)
config = speech.RecognitionConfig(
    encoding=speech.RecognitionConfig.AudioEncoding.LINEAR16,
    sample_rate_hertz=16000,
    language_code="en-US",
)

# Detect speech
response = client.recognize(config=config, audio=audio)

# Print the results
results = []
for result in response.results:
    results.append(result.alternatives[0].transcript)

print("\\n".join(results))
"""
                
                # For the placeholder, simulate a transcription
                transcription = f"Transcription of {metadata.get('filename')}:\n"
                transcription += "This is a placeholder transcription. In a real implementation, "
                transcription += "this would contain the actual text transcribed from the audio file "
                transcription += "using Google Cloud Speech-to-Text API."
                
                return transcription
            finally:
                # Clean up the temporary file
                if os.path.exists(temp_file_path):
                    os.remove(temp_file_path)
                    
        except Exception as e:
            logger.error(f"Error processing audio: {str(e)}")
            return ""


class TextProcessor(BaseProcessor):
    """Processor for text files."""
    
    def process(self, data: bytes, metadata: Dict[str, Any]) -> str:
        """
        Process text data.
        
        Args:
            data: Raw text data
            metadata: Metadata about the file
            
        Returns:
            Processed text
        """
        # Simple decoding with UTF-8 fallback
        try:
            return data.decode('utf-8', errors='replace')
        except Exception as e:
            logger.error(f"Error processing text: {str(e)}")
            return ""


class SharePointProcessor(BaseProcessor):
    """Processor for SharePoint files."""
    
    def process(self, data: bytes, metadata: Dict[str, Any]) -> str:
        """
        Process SharePoint file data.
        
        Args:
            data: Raw file data
            metadata: Metadata about the file
            
        Returns:
            Processed text
        """
        # Determine the file type from metadata and delegate to the appropriate processor
        file_extension = metadata.get('file_extension', '').lower()
        
        if file_extension in ['.pdf']:
            processor = PDFProcessor()
            return processor.process(data, metadata)
        elif file_extension in ['.xlsx', '.xls', '.csv']:
            processor = ExcelProcessor()
            return processor.process(data, metadata)
        elif file_extension in ['.pptx', '.ppt']:
            processor = PowerPointProcessor()
            return processor.process(data, metadata)
        elif file_extension in ['.mp3', '.wav', '.ogg', '.m4a']:
            processor = AudioProcessor()
            return processor.process(data, metadata)
        else:
            # Default to text processor
            processor = TextProcessor()
            return processor.process(data, metadata)


class OneDriveProcessor(SharePointProcessor):
    """Processor for OneDrive files."""
    # Reuse the SharePoint processor since the processing logic is the same
    pass 